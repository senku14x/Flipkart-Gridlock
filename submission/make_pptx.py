"""
submission/make_pptx.py: build a content-rich .pptx pitch deck (accepted Presentation
format). Slide content lives here (not the sparse PITCH_DECK.md outline) so the deck
explains the solution and how it works on its own. Figures are embedded on the
relevant slides; a short presenter note goes in each notes pane.

Run:  pip install python-pptx && python submission/make_pptx.py
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(HERE, "figures")
ACCENT = RGBColor(0xE4, 0x57, 0x2E)
INK = RGBColor(0x1B, 0x2A, 0x4A)
GREY = RGBColor(0x55, 0x5F, 0x70)

SLIDES = [
    {"t": "The problem", "n": "Set up the visibility gap in 20 seconds.", "b": [
        "Bengaluru logs thousands of parking violations a day, but the feed shows where TICKETS are written, not where parking actually blocks traffic.",
        "Those are different places: a lane that catches a daily pre-dawn sweep can out-ticket the junction that jams a main road every evening.",
        "Enforcement is reactive and patrol-based, with no impact heatmap and no way to prioritize.",
        "Goal: tell an officer where, when, and at what cost to deploy, and how many patrols it takes.",
    ]},
    {"t": "The data we have (and don't)", "n": "Be honest about the data up front.", "b": [
        "298,445 real violation records, 151 days (Nov 2023 to Apr 2024), 54 police stations.",
        "Per record: location, violation type, vehicle type, timestamps.",
        "Missing: no traffic speeds, volumes, or travel times, so congestion impact must be DERIVED, not looked up.",
        "Confound: a violation is logged only when a patrol is present. A daily 4-5am sweep is 15% of all records.",
        "So recorded violations are roughly parking demand times patrol presence, and the recorded hour reflects shifts, not congestion.",
    ]},
    {"t": "Our approach: see, understand, act", "fig": "fig-pipeline.png",
     "n": "One pipeline, three layers.", "b": [
        "SEE: an impact-weighted hotspot map and a ranked shortlist of zones.",
        "UNDERSTAND: what the congestion costs, where effort is misallocated, which spots are worsening.",
        "ACT: a forecaster, a false-report triage model, and a patrol optimizer that outputs a costed plan.",
        "Runs on the cleaned data with no live feed, and is built to absorb one when it exists.",
    ]},
    {"t": "How we build it: data foundation", "n": "The unglamorous work that makes the rest valid.", "b": [
        "Convert timestamps UTC to IST. Without this, rush-hour analysis is 5.5 hours wrong.",
        "Map each violation type to an obstruction weight (how much it blocks the carriageway).",
        "Map each vehicle to a PCU footprint: a double-parked bus counts far more than a scooter.",
        "Weight each record by an EXOGENOUS road-utilization curve by hour. This injects traffic flow without a feed and dodges the confound.",
        "Aggregate to H3 res-9 cells (about 150 m): 2,534 hotspot cells. Build the core on approved records.",
    ]},
    {"t": "The Congestion Impact Score (0 to 100)", "fig": "fig-concentration.png",
     "n": "The heart of it. A cell must be bad on several axes to rank high.", "b": [
        "score = 100 x (Volume x Intensity x Exposure x Persistence) ^ 1/4, a geometric mean of percentile ranks.",
        "Volume: how many violations, shrunk so a single-event fluke can't top the list.",
        "Intensity: how obstructive the typical violation is (carriageway, vehicle footprint, junction adjacency).",
        "Exposure: how busy the road normally is. Persistence: how chronic the spot is.",
        "Deliberately decoupled from a raw count (rank correlation 0.56). The top 1% of cells carry 35% of impact.",
    ]},
    {"t": "The impact map: watch the city re-rank", "fig": "fig-map.png",
     "n": "The toggle is the wow moment. Demo it live.", "b": [
        "Every cell (about 150 m) is colored by its Congestion Impact Score.",
        "Toggle raw count to impact: the pre-dawn-sweep clusters fade, the real chokepoints light up.",
        "Hover any cell for its breakdown: the four axes, its estimated cost, and a rising or cooling tag.",
        "The bright core is the known commercial belt: Shivajinagar, Majestic, City Market, Malleshwaram.",
    ]},
    {"t": "Ranked zones: where AND when", "n": "A map is nice; an officer needs a list.", "b": [
        "The top 30 cells, each with its dominant violation, the streets involved, and an exposure-weighted window.",
        "Validated without ground truth: face validity, where 20 of the top 20 are known chokepoints.",
        "Stability: the ranking holds month to month (rank correlation about 0.75 to 0.86).",
        "This is the ops payload: a shortlist an officer can act on tomorrow.",
    ]},
    {"t": "Does it agree with the real city?", "fig": "fig-osm.png",
     "n": "Built blind to OSM, the score still rediscovers the commercial cores. The strongest validation we have.", "b": [
        "The score uses no road network and no land use, just the violation feed. We cross-check it against real OpenStreetMap geography it never saw.",
        "Sort all cells into 10 impact bands: the share next to a market, shop, or transit stop climbs steadily from 34% in the lowest to 62% in the highest.",
        "Top-30 hotspots sit next to a marketplace 17% vs 2% city-wide, an ~8x enrichment; proximity tracks impact about twice as strongly as road class.",
        "Honest read: it barely correlates with arterial road class (0.12), which is correct. Parking chokes narrow commercial streets, not wide arterials.",
    ]},
    {"t": "The visibility gap: effort vs impact", "fig": "fig-gap.png",
     "n": "The persuasion slide. Lead with the night number.", "b": [
        "Recorded violations double as a map of enforcement EFFORT.",
        "Weight each by its impact unit (obstruction x PCU x utilization) to see what that effort actually buys.",
        "The 4-5am sweep: 15% of effort, 4% of impact. The night window: 34% of effort, 8% of impact.",
        "Peak hours are the reverse: 29% of effort, 50% of impact. ParkPulse redirects effort to where it matters.",
    ]},
    {"t": "What it costs", "n": "The number leadership remembers.", "b": [
        "Calibrate the physical delay-potential into vehicle-hours, then rupees, with a low / base / high band.",
        "About 574 vehicle-hours of delay per day; roughly Rs 5.2 crore per year (base case).",
        "The worst 20 cells carry 31% of the cost; just 58 cells carry half.",
        "A first-order estimate, not a measurement. The relative concentration is robust to the assumptions.",
    ]},
    {"t": "Emerging hotspots: early warning", "n": "Proactive, not just descriptive.", "b": [
        "The score ranks where things are bad on average; this flags where they are getting WORSE.",
        "Trend each cell's SHARE of citywide volume over the four full months, which detrends the citywide enforcement ramp.",
        "227 cells are rising; 103 are both high-impact and rising, the priority set.",
        "Honest: a rising trend can reflect shifting enforcement focus, so we call it rising activity.",
    ]},
    {"t": "Forecaster: the one model with real labels", "fig": "fig-forecast.png",
     "n": "This is the genuine ML, and we are honest about its ceiling.", "b": [
        "Predict next-day violation intensity per cell on a strict temporal holdout (train Nov-Feb, test Mar-Apr).",
        "A bake-off across the gradient-boosting family; LightGBM (Tweedie) wins.",
        "Beats the seasonal-naive baseline by 36% on coverage@20; captures 72% of the oracle ceiling.",
        "Honest ablation: 29 extra features and tuning did not beat the base set. The ceiling is the data, not the model.",
    ]},
    {"t": "False-report triage: don't chase ghosts", "n": "A second, cleaner ML win.", "b": [
        "About a third of reviewed reports are rejected on review.",
        "A separate supervised model predicts which, from the report's own attributes (location, time, vehicle, type).",
        "ROC-AUC 0.758 on a stratified holdout: real structure, not a base rate.",
        "Flag the top 20% most-suspect and you catch 43% of all rejections at 64% precision, so patrols go to real problems.",
    ]},
    {"t": "Patrol optimizer + ROI", "fig": "fig-coverage.png",
     "n": "From insight to a deployable, costed plan.", "b": [
        "A greedy max-coverage optimizer assigns N beats (a cell plus its H3 ring) to maximize impact covered.",
        "Because the worst cells cluster, 20 beats cover 53% of citywide impact vs 47% for a naive top-N pick.",
        "The slider also reports the rupee ROI: 20 patrols relieve roughly Rs 77k per day of delay.",
        "Same fleet, more relief. It turns the analysis into a deployable plan with a number attached.",
    ]},
    {"t": "Honest about the gap", "n": "Judges reward candor; so does anyone deploying this.", "b": [
        "No ground truth for impact: an engineered index, validated four ways (face validity, stability, weight-robustness, and an independent OSM cross-check), never claimed as accuracy.",
        "Robust to the weights: across 2,000 random axis weightings the top-20 ranking holds (median rank-corr 0.97), so the score isn't a hand-tuned artefact.",
        "The enforcement confound: we weight by exogenous exposure, never the recorded hour.",
        "The data can't see the evening (enforcement rarely works evenings), so we don't fake an hour-by-hour schedule.",
        "Fusion-ready: a live speed feed turns the index into a learning model, with measured slowdown as the label.",
    ]},
    {"t": "Tech, impact, and what's next", "n": "Close on deployability and the roadmap.", "b": [
        "Pipeline: Python (pandas, h3, scikit-learn, lightgbm). App: Next.js + deck.gl, fully static, on a CDN with no backend or API keys.",
        "Live demo: https://flipkart-gridlock-gamma.vercel.app/",
        "Impact today: a ranked, costed, where-and-when patrol plan with its ROI.",
        "Next: plug in live speeds (a true impact label), the OSM road graph (flow-criticality), and patrol rosters (de-bias the confound).",
    ]},
]


def accent_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank); accent_bar(s, prs)
    tf = textbox(s, 0.9, 2.0, 11.5, 3.6)
    for txt, size, color, bold in [
        ("ParkPulse", 60, ACCENT, True),
        ("Find the parking that actually chokes traffic, and patrol it first.", 24, INK, False),
        ("Gridlock Hackathon 2.0 (Flipkart), Round 2   |   Poor Visibility on Parking-Induced Congestion", 14, GREY, False),
        ("Team Spectres   ·   Kavya Mahajan, Aarav Harshvardhan, Souhardyo Dasgupta, Vishesh Gupta", 16, INK, True),
        ("Demo: https://flipkart-gridlock-gamma.vercel.app/", 14, ACCENT, False),
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = txt; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(10)

    for sl in SLIDES:
        s = prs.slides.add_slide(blank); accent_bar(s, prs)
        has_fig = "fig" in sl and os.path.exists(os.path.join(FIGDIR, sl["fig"]))
        tt = textbox(s, 0.6, 0.5, 12, 1.0)
        p = tt.paragraphs[0]; p.text = sl["t"]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = INK
        body_w = 6.7 if has_fig else 11.9
        size = 15 if len(sl["b"]) >= 5 else 17
        bf = textbox(s, 0.7, 1.65, body_w, 5.4)
        for i, b in enumerate(sl["b"]):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(size); p.font.color.rgb = INK; p.space_after = Pt(10); p.line_spacing = 1.05
        if has_fig:
            s.shapes.add_picture(os.path.join(FIGDIR, sl["fig"]), Inches(7.6), Inches(1.75), width=Inches(5.2))
        if sl.get("n"):
            s.notes_slide.notes_text_frame.text = sl["n"]

    out = os.path.join(HERE, "ParkPulse_Pitch_Deck.pptx")
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
